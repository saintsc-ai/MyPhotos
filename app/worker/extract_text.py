"""Document text extraction for file-domain content search.

Registry keyed by extension. Tier 0 (plain text) is always available; Tier 1
(PDF / office) uses optional libraries — if a library isn't installed the
extractor returns None and the caller records text_status='none' (never a
hard failure, like the OCR opt-in). Extracted text is capped so file_text /
file_fts stay bounded.

extract(abs_path, ext, mime) -> (text | None, engine)
  text=None means "nothing extracted" (unsupported / empty / lib missing);
  engine is a short tag stored in File.text_engine for later reprocessing.
"""

from __future__ import annotations

import logging
from typing import Optional, Tuple

log = logging.getLogger(__name__)

MAX_CHARS = 200_000  # ~plenty for search relevance; keeps the index bounded

# Tier 0 — plain text families (no dependency).
PLAIN_EXTS = {
    "txt", "md", "markdown", "csv", "tsv", "log", "json", "xml", "html", "htm",
    "yaml", "yml", "ini", "cfg", "rst", "tex",
    # common source code
    "py", "js", "ts", "java", "c", "cpp", "h", "hpp", "cs", "go", "rs", "rb",
    "php", "sh", "sql", "css", "kt", "swift",
}


def _cap(s: Optional[str]) -> Optional[str]:
    if not s:
        return None
    s = s.strip()
    if not s:
        return None
    return s[:MAX_CHARS]


def _plain(path: str) -> Optional[str]:
    with open(path, "rb") as fh:
        raw = fh.read(MAX_CHARS * 4)  # bytes; decoded text gets capped below
    for enc in ("utf-8", "cp949", "euc-kr", "latin-1"):
        try:
            return _cap(raw.decode(enc))
        except (UnicodeDecodeError, LookupError):
            continue
    return None


def _pdf(path: str) -> Optional[str]:
    try:
        from pypdf import PdfReader
    except Exception:
        return None
    try:
        reader = PdfReader(path)
        out = []
        total = 0
        for page in reader.pages:
            t = page.extract_text() or ""
            out.append(t)
            total += len(t)
            if total >= MAX_CHARS:
                break
        return _cap("\n".join(out))
    except Exception as e:
        log.warning("pdf extract failed for %s: %s", path, e)
        return None


def _docx(path: str) -> Optional[str]:
    try:
        import docx  # python-docx
    except Exception:
        return None
    try:
        d = docx.Document(path)
        return _cap("\n".join(p.text for p in d.paragraphs))
    except Exception as e:
        log.warning("docx extract failed for %s: %s", path, e)
        return None


def _xlsx(path: str) -> Optional[str]:
    try:
        import openpyxl
    except Exception:
        return None
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        total = 0
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    line = " ".join(cells)
                    parts.append(line)
                    total += len(line)
                    if total >= MAX_CHARS:
                        break
            if total >= MAX_CHARS:
                break
        wb.close()
        return _cap("\n".join(parts))
    except Exception as e:
        log.warning("xlsx extract failed for %s: %s", path, e)
        return None


def _pptx(path: str) -> Optional[str]:
    try:
        from pptx import Presentation
    except Exception:
        return None
    try:
        prs = Presentation(path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
        return _cap("\n".join(parts))
    except Exception as e:
        log.warning("pptx extract failed for %s: %s", path, e)
        return None


def _rtf(path: str) -> Optional[str]:
    try:
        from striprtf.striprtf import rtf_to_text
    except Exception:
        return None
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as fh:
            return _cap(rtf_to_text(fh.read()))
    except Exception as e:
        log.warning("rtf extract failed for %s: %s", path, e)
        return None


def extract(abs_path: str, ext: str, mime: Optional[str]) -> Tuple[Optional[str], str]:
    """Return (text | None, engine). engine records what handled it."""
    e = (ext or "").lower()
    if e in PLAIN_EXTS or (mime or "").startswith("text/"):
        return _plain(abs_path), "plain"
    if e == "pdf":
        return _pdf(abs_path), "pypdf"
    if e == "docx":
        return _docx(abs_path), "python-docx"
    if e == "xlsx":
        return _xlsx(abs_path), "openpyxl"
    if e == "pptx":
        return _pptx(abs_path), "python-pptx"
    if e == "rtf":
        return _rtf(abs_path), "striprtf"
    return None, "unsupported"
